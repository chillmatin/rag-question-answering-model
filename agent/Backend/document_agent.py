import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openai import OpenAI
import re
from collections import deque
import pickle

class DataAgent:
    def __init__(self, api_key, folder_path):
        self.client = OpenAI(api_key=api_key)
        self.folder_path = folder_path
        print("DataAgent working: ")

    def read_pdf(self, file_path):
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    def read_docx(self, file_path):
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def read_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
        
    


    def process_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            content = self.read_pdf(file_path)
        elif ext == '.docx':
            content = self.read_docx(file_path)
        elif ext == '.pptx':
            content = self.read_pptx(file_path)
        return content



    def summarize_content(self, content):
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "Summarize the following content."},
                {"role": "user", "content": f"Summarize: {content}"}
            ],
            max_tokens=300,
            n=1,
            temperature=0.5
        )
        return response.choices[0].message.content.strip()

    def extract_keywords(self, content):
        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": 
                """
                Your task is to extract keywords from a prompt to do an algorithmic (non llm) word search. Answer should be a string that has 15 keywords that will help us find the asked question from the data at hand.
                """},
                {"role": "user", "content":  f"Extract from this text: {content}" }
            ],
            max_tokens=32
        )
        return response.choices[0].message.content.strip()


    def get_answer(self, prompt):        
        
        keywords = self.extract_keywords(prompt)
        context = self.extract_from_files(self.file_summaries,keywords)
        evaluation = "\n".join(context)

        answer = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": 
                f"""
                Your task is to give a precise answer to the user prompt based on the information given on the {evaluation}.
                """},
                {"role": "user", "content":  f"Answer from this text: {prompt}" }
            ],
            max_tokens=150,
            temperature=0.5,
            n=1
        )


        return answer.choices[0].message.content.strip()
    
    
    def extract_context(self, corpus, target_words, window_size=100, save_positions=True, positions_file='positions.pkl'):
        words = corpus.split()
        lower_words = [word.lower() for word in words]
        target_words = [word.lower() for word in target_words]
        
        if save_positions:
            try:
                with open(positions_file, 'rb') as file:
                    word_positions = pickle.load(file)
            except FileNotFoundError:
                word_positions = {word: [] for word in target_words}
                for index, word in enumerate(lower_words):
                    if word in word_positions:
                        word_positions[word].append(index)
                with open(positions_file, 'wb') as file:
                    pickle.dump(word_positions, file)
        else:
            word_positions = {word: [] for word in target_words}
            for index, word in enumerate(lower_words):
                if word in word_positions:
                    word_positions[word].append(index)

        contexts = []
        combined_context = []
        last_end = -1
        
        for word, positions in word_positions.items():
            for pos in positions:
                start = max(pos - window_size, 0)
                end = min(pos + window_size + 1, len(words))
                
                if start <= last_end:
                    combined_context.extend(words[last_end:end])
                else:
                    if combined_context:
                        contexts.append(' '.join(combined_context))
                    combined_context = words[start:end]
                    
                last_end = end
        
        if combined_context:
            contexts.append(' '.join(combined_context))
        
        return contexts

    def extract_from_files(self, file_summaries, target_words, window_size=50):
        corpus = ""
        for file_path in file_summaries.keys():
            if file_path.endswith('.pdf'):
                corpus += self.read_pdf(file_path) + "\n"
            elif file_path.endswith('.docx'):
                corpus += self.read_docx(file_path) + "\n"
            elif file_path.endswith('.pptx'):
                corpus += self.read_pptx(file_path) + "\n"

        contexts = self.extract_context(corpus, target_words, window_size, save_positions=True, positions_file="positions.pkl")
        return contexts
    


# # # Example usage:
# agent = DataAgent(api_key='key here', folder_path='/home/cuneyd/agent/Backend/files')
# prompt = "I fought with my coworker, what is going to happen to me?"
# answer = agent.get_answer(prompt)
# print(answer)
